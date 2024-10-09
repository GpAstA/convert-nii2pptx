# -*- coding: utf-8 -*-
"""
Created on Mon Jul  4 13:11:29 2022

@author: ME-pc
"""

from pptx import Presentation
from pptx.util import Inches, Pt
import datetime

from matplotlib import pylab as plt
import nibabel as nib

import numpy as np
from PIL import Image
import os
import imageio

from tkinter import filedialog

SLIDE_WIDTH, SLIDE_HEIGHT = 9144000, 6858000
IMG_CENTER_X, IMG_CENTER_Y = SLIDE_WIDTH / 2, SLIDE_HEIGHT / 2
SLIDE_ASPECT_RATIO = SLIDE_WIDTH / SLIDE_HEIGHT


def add_slide(prs):
    # 白紙スライドの追加(ID=6は白紙スライド)
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    return slide


def add_picture(slide, save_png, i):
    # 画像サイズを取得してアスペクト比を得る
    im = Image.open(save_png)
    im_width, im_height = im.size
    aspect_ratio = im_width / im_height

    # スライドと画像のアスペクト比に応じて処理を分岐
    if aspect_ratio > SLIDE_ASPECT_RATIO:
        img_display_width = SLIDE_WIDTH
        img_display_height = img_display_width / aspect_ratio
    else:
        img_display_height = SLIDE_HEIGHT
        img_display_width = img_display_height * aspect_ratio

    # 画像の位置をセンタリング
    left = IMG_CENTER_X - img_display_width / 2
    top = IMG_CENTER_Y - img_display_height / 2
    left2 = IMG_CENTER_X + img_display_width / 2 + 10
    top2 = IMG_CENTER_X - img_display_width / 2 + 10

    # 画像とテキストを追加
    slide.shapes.add_picture(save_png, left, top, width=img_display_width if aspect_ratio > SLIDE_ASPECT_RATIO else None, height=None if aspect_ratio > SLIDE_ASPECT_RATIO else img_display_height)
    txbox = slide.shapes.add_textbox(left2, top2, Inches(1), Inches(1))
    tf = txbox.text_frame
    tf.text = 'z={}'.format(i)
    tf.paragraphs[0].font.size = Pt(28)
    return slide


def create_pptx_for_nii(nii_file, output_dir):
    now = str(datetime.datetime.now().time())

    # 出力フォルダを作成
    os.makedirs(output_dir, exist_ok=True)
    
    # ファイル名からpptxの保存名を設定
    base_filename = os.path.basename(nii_file).replace('.nii', '.pptx')
    save_name = os.path.join(output_dir, base_filename)

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    img = nib.load(nii_file)
    img_fdata = img.get_fdata()
    z_slices = img.shape[2]

    for i in range(z_slices):
        slice_data = img_fdata[:, :, i]
        slice_data = np.rot90(slice_data)
        save_png = os.path.join(output_dir, f'z={i}.png')
        imageio.imwrite(save_png, slice_data)

        slide = add_slide(prs)
        add_picture(slide, save_png, i)

        # 一時的なPNGを削除
        os.remove(save_png)

    prs.save(save_name)
    print(f"ファイルの書き出し完了しました: {save_name}")


def find_nii_and_create_pptx(input_dir, output_base_dir):
    for root, dirs, files in os.walk(input_dir):
        for file in files:
            if file.endswith(".nii"):
                nii_file = os.path.join(root, file)
                
                # 出力先のディレクトリ構造を維持
                relative_path = os.path.relpath(root, input_dir)
                output_dir = os.path.join(output_base_dir, relative_path)
                
                # PPTX作成
                create_pptx_for_nii(nii_file, output_dir)


# メイン処理
if __name__ == "__main__":
    typ = [('Nii files', '*.nii')] 
    input_dir = filedialog.askdirectory(title="フォルダを選択してください")
    output_dir = filedialog.askdirectory(title="出力フォルダを選択してください")

    find_nii_and_create_pptx(input_dir, output_dir)
