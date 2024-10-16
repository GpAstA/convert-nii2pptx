import os
from pptx import Presentation
from tkinter import Tk, filedialog

# ファイル選択ダイアログを表示する関数
def select_file(title):
    root = Tk()
    root.withdraw()  # メインウィンドウを表示しない
    file_path = filedialog.askopenfilename(title=title, filetypes=[("PowerPoint files", "*.pptx")])
    return file_path

# フォルダを選択する関数（保存先フォルダ）
def select_folder(title):
    root = Tk()
    root.withdraw()
    folder_path = filedialog.askdirectory(title=title)
    return folder_path

# シェイプの描画オブジェクトを抽出する再帰関数
def extract_shapes(slide, new_slide):
    for shape in slide.shapes:
        if hasattr(shape, 'shapes') and len(shape.shapes) > 0:  # GroupShape の場合
            extract_shapes(shape, new_slide)  # グループ内のシェイプを再帰的に処理
        elif shape.shape_type in [6, 3]:  # 6 = line, 3 = freeform
            # 位置とサイズを取得
            left = shape.left
            top = shape.top
            width = shape.width
            height = shape.height

            # フリーフォームや線のシェイプのみ追加
            if shape.shape_type == 6:  # 線
                new_shape = new_slide.shapes.add_shape(
                    shape.auto_shape_type, left, top, width, height
                )
                new_shape.line.color.rgb = shape.line.color.rgb
                new_shape.line.width = shape.line.width
            elif shape.shape_type == 3:  # フリーフォーム
                # 新しいスライドにフリーフォームを追加
                new_freeform = new_slide.shapes.add_freeform_shape(
                    shape.left, shape.top, shape.width, shape.height
                )
                new_freeform.fill.fore_color.rgb = shape.fill.fore_color.rgb

# 手書きの線（line, freeform）だけを抽出する関数
def extract_handwriting(input_file, output_folder):
    # 元のPowerPointファイルを読み込み
    presentation = Presentation(input_file)
    
    # 新しいプレゼンテーションを作成
    new_presentation = Presentation()

    # 元のスライドから手書きの線（シェイプ）を抽出
    for slide in presentation.slides:
        new_slide_layout = new_presentation.slide_layouts[5]  # 空白スライドレイアウト
        new_slide = new_presentation.slides.add_slide(new_slide_layout)

        extract_shapes(slide, new_slide)  # シェイプを再帰的に抽出

    # 元のファイル名に"_draw"を付加して新しいファイル名を作成
    base_filename = os.path.splitext(os.path.basename(input_file))[0]  # 元ファイル名（拡張子除去）
    output_filename = f"{base_filename}_draw.pptx"  # 新しいファイル名
    output_path = os.path.join(output_folder, output_filename)

    # 新しいPowerPointファイルとして保存
    new_presentation.save(output_path)
    print(f"手書きの部分だけを含む新しいファイルが保存されました: {output_path}")

# メイン処理
if __name__ == "__main__":
    # 入力ファイルを選択
    input_file = select_file("抽出元のPowerPointファイルを選択してください")
    if not input_file:
        print("ファイルが選択されませんでした。処理を終了します。")
    else:
        # 保存先フォルダを選択
        output_folder = select_folder("保存するフォルダを選択してください")
        if not output_folder:
            print("保存先フォルダが指定されませんでした。処理を終了します。")
        else:
            # 手書き部分を抽出して新しいファイルを作成
            extract_handwriting(input_file, output_folder)
