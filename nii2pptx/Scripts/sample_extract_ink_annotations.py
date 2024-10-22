from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
import numpy as np
import nibabel as nib
import os
from tkinter import Tk, filedialog

def is_target_color(color):
    """指定されたペンの色と一致するか判定する関数"""
    target_colors = [
        (51, 204, 255),   # 青
        (231, 18, 36),    # 赤
        (102, 204, 0),    # 緑
        (255, 193, 20)    # 黄
    ]
    return (color[0], color[1], color[2]) in target_colors

def process_ppt(ppt_file, output_dir):
    ppt = Presentation(ppt_file)

    # 各スライドを処理
    for slide_num, slide in enumerate(ppt.slides):
        # スライドごとに画像として保存するためのキャンバス作成
        slide_width = ppt.slide_width
        slide_height = ppt.slide_height
        scale_factor = 0.001  # ここで縮小率を指定
        new_width = int(slide_width * scale_factor)
        new_height = int(slide_height * scale_factor)
        # image = Image.new('RGB', (slide_width, slide_height), (255, 255, 255))
        image = Image.new('RGB', (new_width, new_height), (255, 255, 255))

        draw = ImageDraw.Draw(image)

        for shape in slide.shapes:
            # # テキストを削除
            # if shape.has_text_frame:
            #     continue  # テキスト要素を無視

            # # 中央の白黒画像（9144000x6858000）を削除
            # if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            #     if shape.width == 9144000 and shape.height == 6858000:
            #         continue  # 背景画像を無視

            # 手書きのペンの線で、指定された色のみを保持
            if shape.shape_type == MSO_SHAPE_TYPE.FREEFORM:
                if hasattr(shape, "line") and hasattr(shape.line, "fill") and hasattr(shape.line.fill, "fore_color"):
                    color = shape.line.fill.fore_color.rgb
                    if is_target_color(color):  # 指定されたペンの色と一致する場合のみ処理
                        for point in shape.points:
                            draw.point((point.x, point.y), fill=(color[0], color[1], color[2]))

        # スライドを512x512にリサイズ
        resized_image = resize_and_crop_to_512(image)
        
        # スライド画像を保存
        slide_image_path = os.path.join(output_dir, f"slide_{slide_num}.png")
        resized_image.save(slide_image_path)

        # ここで、手書き部分の画像をNIfTIに変換する
        convert_image_to_nii(slide_image_path, os.path.join(output_dir, f"slide_{slide_num}.nii"))

# スライドを512x512にリサイズし、中央を切り取る関数
def resize_and_crop_to_512(image):
    # 短い辺を512ピクセルにリサイズ
    width, height = image.size
    if width > height:
        new_height = 512
        new_width = int((width / height) * 512)
    else:
        new_width = 512
        new_height = int((height / width) * 512)

    resized_image = image.resize((new_width, new_height))

    # 中央部分を512x512にクロップ
    left = (new_width - 512) / 2
    top = (new_height - 512) / 2
    right = left + 512
    bottom = top + 512

    return resized_image.crop((left, top, right, bottom))

# 画像をNIfTI形式に変換する関数
def convert_image_to_nii(image_path, output_filename):
    img = Image.open(image_path).convert('L')  # グレースケール変換
    image_data = np.array(img)
    nii_image = nib.Nifti1Image(image_data, np.eye(4))
    nib.save(nii_image, output_filename)

# tkを使用してファイルとフォルダを選択
def select_file_and_folder():
    root = Tk()
    root.withdraw()  # メインウィンドウを隠す

    # PowerPointファイルの選択
    ppt_file = filedialog.askopenfilename(title="Select PowerPoint file", filetypes=[("PowerPoint files", "*.pptx")])
    if not ppt_file:
        print("No PowerPoint file selected. Exiting.")
        return None, None

    # 出力フォルダの選択
    output_dir = filedialog.askdirectory(title="Select output folder")
    if not output_dir:
        print("No output directory selected. Exiting.")
        return None, None

    return ppt_file, output_dir

if __name__ == "__main__":
    # PowerPointファイルと出力フォルダを選択
    ppt_file, output_dir = select_file_and_folder()

    if ppt_file and output_dir:
        # 出力フォルダが存在しない場合は作成
        if not os.path.exists(output_dir):
            os.makedirs(output_dir)

        # PowerPointファイルを処理
        process_ppt(ppt_file, output_dir)
