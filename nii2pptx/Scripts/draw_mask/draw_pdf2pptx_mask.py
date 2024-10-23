import os
import fitz  # PyMuPDF
from PIL import Image
import numpy as np
import cv2
import statistics
from copy import deepcopy
from pptx import Presentation
from pptx.util import Inches
import io
import tkinter as tk
from tkinter import filedialog

# PDFから画像を抽出する関数
def get_pdf_images(pdf_path):
    doc = fitz.open(pdf_path)
    images_list = []

    for page_num in range(doc.page_count):
        page = doc[page_num]
        img = page.get_pixmap()

        # Convert the RGB image to a PIL Image
        pil_img = Image.frombytes("RGB", (img.width, img.height), img.samples)

        # Append the PIL Image to the list
        images_list.append(pil_img)

    doc.close()
    return images_list

# PIL画像をNumPy配列に変換
def images_to_numpy(images_list):
    return [np.array(img) for img in images_list]

# マスク抽出（緑色の範囲を広げて調整）
def mask_extracted(image_array):
    combined_images = []
    for i in range(len(image_array)):
        image_rgb = image_array[i]

        # 赤い部分の抽出
        red_mask = cv2.inRange(image_rgb, np.array([150, 0, 0]), np.array([255, 100, 100]))
        red_extracted = cv2.bitwise_and(image_rgb, image_rgb, mask=red_mask)

        # 水色の部分の抽出
        cyan_mask = cv2.inRange(image_rgb, np.array([0, 100, 150]), np.array([100, 255, 255]))
        cyan_extracted = cv2.bitwise_and(image_rgb, image_rgb, mask=cyan_mask)

        # オレンジの部分の抽出
        orange_mask = cv2.inRange(image_rgb, np.array([200, 100, 0]), np.array([255, 200, 100]))
        orange_extracted = cv2.bitwise_and(image_rgb, image_rgb, mask=orange_mask)

        # 黄緑の部分の抽出（範囲を広げたバージョン）
        lower_green = np.array([50, 100, 0])  # 最小値
        upper_green = np.array([130, 255, 50])  # 最大値
        green_mask = cv2.inRange(image_rgb, lower_green, upper_green)
        green_extracted = cv2.bitwise_and(image_rgb, image_rgb, mask=green_mask)

        # 赤、水色、オレンジ、黄緑の部分を合成
        combined_image = cv2.add(red_extracted, cv2.add(cyan_extracted, cv2.add(orange_extracted, green_extracted)))
        combined_images.append(combined_image)

    return combined_images

# リサイズとトリミング
def resize_and_crop(images):
    cropped_images = []
    for image in images:
        height, width = image.shape[:2]

        # 左右中心から切り取る範囲を計算
        crop_start = (width - height) // 2
        crop_end = crop_start + height

        # 切り取り
        cropped_image = image[:, crop_start:crop_end, :]
        cropped_images.append(cropped_image)

    return cropped_images

# 領域のマスキング処理
def masking(cropped_images):
    mask_images = []
    for cropped_image in cropped_images:
        img = cropped_image
        img = cv2.cvtColor(img, cv2.COLOR_BGR2RGB)

        # 最頻値を調べる
        R_mode = statistics.mode(img[:, :, 0].flatten())
        G_mode = statistics.mode(img[:, :, 1].flatten())
        B_mode = statistics.mode(img[:, :, 2].flatten())

        # 領域内を塗りつぶすためのマスクを作成
        mask = np.zeros((img.shape[0], img.shape[1]), dtype=np.uint8)
        mask[(img[:, :, 0] != R_mode) & (img[:, :, 1] != G_mode) & (img[:, :, 2] != B_mode)] = 255

        # 輪郭の検出
        contours, hierarchy = cv2.findContours(mask, cv2.RETR_LIST, cv2.CHAIN_APPROX_SIMPLE)

        # 塗りつぶし処理（白で塗りつぶす）
        img_with_area = deepcopy(img)
        for i in range(len(contours)):
            cv2.fillPoly(img_with_area, [contours[i]], (0, 255, 0))  # 白で内部を塗りつぶし

        mask_images.append(img_with_area)

    return mask_images

# PowerPointに画像を保存する関数
def save_images_to_pptx(images, pptx_path):
    # PowerPointプレゼンテーションを作成
    prs = Presentation()

    # スライドの幅と高さを取得
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    for image in images:
        # スライドを追加
        slide_layout = prs.slide_layouts[6]  # 空白のスライド
        slide = prs.slides.add_slide(slide_layout)

        # 画像をPILからバイトIOに変換してpptxに挿入
        img_bytes = io.BytesIO()
        img_pil = Image.fromarray(image)
        img_pil.save(img_bytes, format='PNG')
        img_bytes.seek(0)

        # 画像の元のサイズを取得
        img_width, img_height = img_pil.size

        # スライドのアスペクト比に合わせてリサイズ（幅と高さのどちらかに合わせる）
        width_ratio = slide_width / img_width
        height_ratio = slide_height / img_height
        scale_factor = min(width_ratio, height_ratio)  # スライドに収まるようにリサイズ

        # 新しい画像サイズ（アスペクト比を保つ）
        new_img_width = int(img_width * scale_factor)
        new_img_height = int(img_height * scale_factor)

        # 画像を中央に配置するための座標計算
        x = (slide_width - new_img_width) // 2  # 中央配置のためのX座標
        y = (slide_height - new_img_height) // 2  # 中央配置のためのY座標

        # スライドに画像を追加（位置とサイズを指定）
        slide.shapes.add_picture(img_bytes, x, y, width=new_img_width, height=new_img_height)

    # pptxを保存
    prs.save(pptx_path)

# PDFファイル名から出力ファイル名を生成する関数
def generate_output_pptx_path(pdf_path):
    pdf_filename = os.path.basename(pdf_path)  # 元のPDFファイル名を取得
    pptx_filename = f"auto_drawed_{os.path.splitext(pdf_filename)[0]}.pptx"  # 出力ファイル名を作成
    return os.path.join(os.path.dirname(pdf_path), pptx_filename)

# 指定されたフォルダ内の全PDFに対して処理を実行
def process_folder(input_folder):
    # フォルダ内の全ファイルを探索
    for root, dirs, files in os.walk(input_folder):
        for file in files:
            if file.endswith(".pdf"):
                pdf_path = os.path.join(root, file)
                print(f"処理中: {pdf_path}")

                # PDFから画像を取得
                pdf_images = get_pdf_images(pdf_path)

                # Convert PIL Images to NumPy arrays
                numpy_images = images_to_numpy(pdf_images)

                # マスク処理
                combined_images = mask_extracted(numpy_images)

                # リサイズおよびトリミング
                cropped_images = resize_and_crop(combined_images)

                # マスキング処理
                mask_images = masking(cropped_images)

                # 出力PPTXファイルパスを生成
                pptx_output_path = generate_output_pptx_path(pdf_path)

                # 画像をPPTXに保存
                save_images_to_pptx(mask_images, pptx_output_path)

                print(f"保存完了: {pptx_output_path}")

# tkinterでフォルダ選択ダイアログを表示
def select_folder():
    root = tk.Tk()
    root.withdraw()  # メインウィンドウを隠す
    folder_selected = filedialog.askdirectory()  # フォルダ選択ダイアログを開く
    return folder_selected

if __name__ == "__main__":
    # フォルダ選択ダイアログを開いて、処理するフォルダを選択
    input_folder = select_folder()

    if input_folder:
        process_folder(input_folder)
    else:
        print("フォルダが選択されていません。")
