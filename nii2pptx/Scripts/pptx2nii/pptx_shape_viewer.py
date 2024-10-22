from pptx import Presentation
from tkinter import Tk, filedialog

# ファイル選択ダイアログを表示する関数
def select_file(title):
    root = Tk()
    root.withdraw()  # メインウィンドウを表示しない
    file_path = filedialog.askopenfilename(title=title, filetypes=[("PowerPoint files", "*.pptx")])
    return file_path

# シェイプの種類を文字列で取得
def get_shape_type(shape):
    shape_type_dict = {
        1: "Text Box",
        2: "Table",
        3: "Freeform (Handwritten or Drawn)",
        4: "Picture",
        5: "Media",
        6: "Line",
        7: "Group",
        8: "SmartArt",
        9: "Chart",
        10: "Shape"
    }
    return shape_type_dict.get(shape.shape_type, "Unknown")

# スライドのシェイプ情報を出力
def print_slide_shapes(presentation):
    for slide_index, slide in enumerate(presentation.slides):
        print(f"\n--- Slide {slide_index + 1} ---")
        for shape_index, shape in enumerate(slide.shapes):
            shape_type = get_shape_type(shape)
            print(f"  Shape {shape_index + 1}:")
            print(f"    Type: {shape_type}")
            print(f"    Position: (left: {shape.left}, top: {shape.top})")
            print(f"    Size: (width: {shape.width}, height: {shape.height})")

            # テキストボックスの場合はテキスト内容を表示
            if shape.has_text_frame:
                print(f"    Text: {shape.text}")

            # グループ化されたシェイプの場合、再帰的に中身を表示
            if shape_type == "Group":
                for sub_shape_index, sub_shape in enumerate(shape.shapes):
                    sub_shape_type = get_shape_type(sub_shape)
                    print(f"      Sub-shape {sub_shape_index + 1}:")
                    print(f"        Type: {sub_shape_type}")
                    print(f"        Position: (left: {sub_shape.left}, top: {sub_shape.top})")
                    print(f"        Size: (width: {sub_shape.width}, height: {sub_shape.height})")
                    if sub_shape.has_text_frame:
                        print(f"        Text: {sub_shape.text}")

# メイン処理
if __name__ == "__main__":
    # 入力ファイルを選択
    input_file = select_file("PowerPointファイルを選択してください")
    if not input_file:
        print("ファイルが選択されませんでした。処理を終了します。")
    else:
        # PowerPointファイルの読み込み
        presentation = Presentation(input_file)
        
        # 各スライドのシェイプ情報を出力
        print_slide_shapes(presentation)
