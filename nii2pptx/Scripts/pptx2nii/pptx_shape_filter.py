import os
from pptx import Presentation
from tkinter import Tk, filedialog

# ファイル選択ダイアログを表示する関数
def select_file(title):
    root = Tk()
    root.withdraw()  # メインウィンドウを表示しない
    file_path = filedialog.askopenfilename(title=title, filetypes=[("PowerPoint files", "*.pptx")])
    return file_path

# フォルダを選択する関数（保存先フォルダ）
def select_folder(title):
    root = Tk()
    root.withdraw()
    folder_path = filedialog.askdirectory(title=title)
    return folder_path

# シェイプの抽出処理
def extract_and_save_filtered_presentation(input_file, output_folder):
    # 元のPowerPointファイルを読み込み
    presentation = Presentation(input_file)

    # 新しいプレゼンテーションを作成
    new_presentation = Presentation()

    # スライドごとに処理
    for slide in presentation.slides:
        # 新しいプレゼンテーションに空白のスライドを追加
        new_slide_layout = new_presentation.slide_layouts[5]  # 空白スライド
        new_slide = new_presentation.slides.add_slide(new_slide_layout)

        # 各シェイプを処理（Shape 1およびテキストを含むShapeを除外）
        for shape in slide.shapes:
            # Shape 1（大きな画像）を削除（位置やサイズで判別）
            if shape.left == 1143000 and shape.width == 6858000:
                continue  # このシェイプは除外

            # テキストが含まれるシェイプを削除
            if shape.has_text_frame and shape.text.startswith("z="):
                continue  # このシェイプは除外

            # 他のシェイプ（手書きの線など）を新しいスライドにコピー
            new_shape = new_slide.shapes.add_shape(
                shape.auto_shape_type, shape.left, shape.top, shape.width, shape.height
            )

            # シェイプのフォーマット（線の色など）をコピー
            new_shape.fill.fore_color.rgb = shape.fill.fore_color.rgb
            new_shape.line.color.rgb = shape.line.color.rgb
            new_shape.line.width = shape.line.width

    # 保存ファイル名を元のファイル名に_try_drawを付けて生成
    base_filename = os.path.splitext(os.path.basename(input_file))[0]  # 元ファイル名（拡張子除去）
    output_filename = f"{base_filename}_try_draw.pptx"  # 新しいファイル名
    output_path = os.path.join(output_folder, output_filename)

    # 新しいPowerPointファイルとして保存
    new_presentation.save(output_path)
    print(f"不要なシェイプを削除した新しいファイルが保存されました: {output_path}")

# メイン処理
if __name__ == "__main__":
    # 入力ファイルを選択
    input_file = select_file("抽出元のPowerPointファイルを選択してください")
    if not input_file:
        print("ファイルが選択されませんでした。処理を終了します。")
    else:
        # 保存先フォルダを選択
        output_folder = select_folder("保存するフォルダを選択してください")
        if not output_folder:
            print("保存先フォルダが指定されませんでした。処理を終了します。")
        else:
            # 不要なシェイプを削除して新しいファイルを作成
            extract_and_save_filtered_presentation(input_file, output_folder)
